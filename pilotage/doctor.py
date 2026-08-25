"""Read-only deployment readiness checks.

Hermes's doctor pattern is kept: every probe is failure-isolated and the full
report is shown even when early checks fail.  Pilotage narrows the probes to
its designated production contract and returns a failing exit code whenever a
required capability is not ready.
"""

from __future__ import annotations

import asyncio
import importlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Awaitable, Callable, Optional

import httpx

from .agent import Agent
from .channels.whatsapp import (
    WhatsAppSessionError,
    normalize_whatsapp_chat_id,
    validate_whatsapp_session,
)
from .channels.telegram import normalize_telegram_home_chat_id
from .codex import auth
from .history import ConversationStore
from .redact import redact_sensitive_text
from .runtime_lock import RuntimeLockError, runtime_lock_is_held
from .tools import ToolContext, build_registry, enabled_groups
from .tools.code_execution import _execute
from .tools.subprocess_env import PROTECTED_ENV_VARS, build_subprocess_env


_MAIN_IMPORTS = (
    ("openai", "OpenAI SDK"),
    ("PIL", "Pillow"),
    ("httpx", "HTTPX"),
    ("telegram", "Telegram"),
    ("yaml", "PyYAML"),
    ("croniter", "croniter"),
    ("ddgs", "DDGS"),
    ("firecrawl", "Firecrawl"),
)
_SYSTEM_COMMANDS = (
    "ffmpeg",
    "ffprobe",
    "fc-match",
    "libreoffice",
    "pandoc",
    "pdftotext",
    "pdftoppm",
    "qpdf",
    "rg",
    "tesseract",
)
_CHROMIUM_CANDIDATES = (
    "chromium",
    "chromium-browser",
    "google-chrome",
    "google-chrome-stable",
)
_SQL_ENV_NAMES = (
    "MSSQL_HOST",
    "MSSQL_USER",
    "MSSQL_PASSWORD",
    "MSSQL_DB",
)

_CHART_SMOKE = """
import matplotlib
import numpy

matplotlib.use("Agg")
import matplotlib.pyplot as plt

figure, axis = plt.subplots()
axis.plot(numpy.array([1, 2, 3]), numpy.array([1, 4, 9]))
figure.savefig("chart.png")
plt.close(figure)
"""

_DOCS_SMOKE = """
from pathlib import Path

import defusedxml
import lxml
import pandas
import pdfplumber
import pytesseract
import reportlab
from docx import Document
from markitdown import MarkItDown
from openpyxl import Workbook
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

root = Path(".")
document = Document()
document.add_paragraph("Pilotage")
document.save(root / "document.docx")
assert Document(root / "document.docx").paragraphs[0].text == "Pilotage"

workbook = Workbook()
workbook.active["A1"] = "Pilotage"
workbook.save(root / "document-workbook.xlsx")

presentation = Presentation()
presentation.slides.add_slide(presentation.slide_layouts[6])
presentation.save(root / "presentation.pptx")

writer = PdfWriter()
writer.add_blank_page(width=595, height=842)
writer.write(root / "source.pdf")
with pdfplumber.open(root / "source.pdf") as pdf:
    assert len(pdf.pages) == 1
assert convert_from_path(root / "source.pdf", first_page=1, last_page=1)
assert MarkItDown()
assert pytesseract.get_tesseract_version()
assert reportlab.Version
assert lxml.__version__
assert defusedxml

Image.new("RGB", (16, 16), "white").save(root / "image.png")
pandas.DataFrame({"value": [1]}).to_csv(root / "data.csv", index=False)
"""

_EXCEL_SMOKE = """
import xlsxwriter

workbook = xlsxwriter.Workbook("workbook.xlsx")
worksheet = workbook.add_worksheet("Verification")
worksheet.write(0, 0, "Pilotage")
workbook.close()
"""

_PDF_SMOKE = """
from weasyprint import HTML

html = '''
<html lang="ar" dir="rtl">
<meta charset="utf-8">
<style>
body { font-family: "Noto Sans Arabic", "Noto Sans", sans-serif; }
</style>
<body><h1>جاهز</h1><p lang="fr" dir="ltr">Pilotage est prêt.</p></body>
</html>
'''
HTML(string=html).write_pdf("bilingual.pdf")
"""

_SMOKE_ARTIFACTS = {
    "chart": ("chart.png",),
    "docs": (
        "document.docx",
        "document-workbook.xlsx",
        "presentation.pptx",
        "source.pdf",
        "image.png",
        "data.csv",
    ),
    "excel": ("workbook.xlsx",),
    "pdf": ("bilingual.pdf",),
}
_SMOKE_CODE = {
    "chart": _CHART_SMOKE,
    "docs": _DOCS_SMOKE,
    "excel": _EXCEL_SMOKE,
    "pdf": _PDF_SMOKE,
}


class DoctorError(RuntimeError):
    """One readiness probe failed."""


@dataclass(frozen=True)
class CheckResult:
    name: str
    ok: bool
    detail: str = ""


@dataclass
class DoctorReport:
    profile: str
    checks: list[CheckResult] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return bool(self.checks) and all(check.ok for check in self.checks)

    @property
    def failures(self) -> list[CheckResult]:
        return [check for check in self.checks if not check.ok]


def _safe_detail(value: Any) -> str:
    text = str(value or "")
    for name in (
        *PROTECTED_ENV_VARS,
        "MSSQL_PASSWORD",
        "FIRECRAWL_API_KEY",
    ):
        secret = os.environ.get(name, "")
        if secret:
            text = text.replace(secret, "[REDACTED]")
    text = redact_sensitive_text(text).strip()
    return text[:700]


async def _probe(
    report: DoctorReport,
    name: str,
    function: Callable[[], Any],
) -> None:
    try:
        detail = await asyncio.to_thread(function)
    except Exception as exc:
        report.checks.append(
            CheckResult(name, False, _safe_detail(exc))
        )
    else:
        report.checks.append(
            CheckResult(name, True, _safe_detail(detail))
        )


async def _probe_async(
    report: DoctorReport,
    name: str,
    function: Callable[[], Awaitable[Any]],
) -> None:
    try:
        detail = await function()
    except Exception as exc:
        report.checks.append(
            CheckResult(name, False, _safe_detail(exc))
        )
    else:
        report.checks.append(
            CheckResult(name, True, _safe_detail(detail))
        )


def _run(
    command: list[str],
    *,
    cwd: Optional[Path] = None,
    timeout: float = 20,
    env: Optional[dict[str, str]] = None,
) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        command,
        cwd=str(cwd) if cwd is not None else None,
        env=env or build_subprocess_env(),
        stdin=subprocess.DEVNULL,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=timeout,
        check=False,
        creationflags=(
            int(getattr(subprocess, "CREATE_NO_WINDOW", 0))
            if os.name == "nt"
            else 0
        ),
    )


def _check_main_python() -> str:
    if not ((3, 11) <= sys.version_info < (3, 14)):
        raise DoctorError(
            f"Python 3.11 through 3.13 is required, found {sys.version.split()[0]}"
        )
    missing: list[str] = []
    for module, label in _MAIN_IMPORTS:
        try:
            importlib.import_module(module)
        except Exception:
            missing.append(label)
    if missing:
        raise DoctorError("missing imports: " + ", ".join(missing))
    return f"Python {sys.version.split()[0]}, {len(_MAIN_IMPORTS)} imports"


def _check_command(name: str) -> str:
    path = shutil.which(name)
    if not path:
        raise DoctorError(f"{name} is not on PATH")
    return path


def _check_node() -> str:
    node = shutil.which("node")
    npm = shutil.which("npm")
    if not node or not npm:
        raise DoctorError("Node.js and npm are required")
    result = _run([node, "-p", "process.versions.node"])
    version = result.stdout.strip()
    match = re.match(r"(\d+)", version)
    if result.returncode != 0 or match is None:
        raise DoctorError("could not read the Node.js version")
    if int(match.group(1)) < 20:
        raise DoctorError(f"Node.js 20+ is required, found {version}")
    return f"Node.js {version}"


def _check_bridge_dependencies(config: Any) -> str:
    bridge_dir = Path(config.bridge_dir)
    if not config.bridge_script.is_file():
        raise DoctorError(f"bridge script is missing: {config.bridge_script}")
    node = shutil.which("node")
    if not node:
        raise DoctorError("node is not on PATH")
    code = (
        "await import('@whiskeysockets/baileys');"
        "await import('express');"
        "await import('pino');"
        "await import('qrcode-terminal');"
    )
    result = _run(
        [node, "--input-type=module", "-e", code],
        cwd=bridge_dir,
        timeout=30,
    )
    if result.returncode != 0:
        raise DoctorError(
            "bridge Node dependencies are not loadable: "
            + _safe_detail(result.stderr)
        )
    return "bridge modules load"


def _check_tesseract_languages() -> str:
    executable = _check_command("tesseract")
    result = _run([executable, "--list-langs"])
    if result.returncode != 0:
        raise DoctorError("tesseract --list-langs failed")
    languages = {
        line.strip()
        for line in (result.stdout + "\n" + result.stderr).splitlines()
        if line.strip()
    }
    missing = sorted({"ara", "eng", "fra"} - languages)
    if missing:
        raise DoctorError(
            "missing Tesseract languages: " + ", ".join(missing)
        )
    return "ara, eng, fra"


def _check_arabic_font() -> str:
    executable = _check_command("fc-match")
    result = _run(
        [executable, "-f", "%{family}", "Noto Sans Arabic"]
    )
    family = result.stdout.strip()
    if result.returncode != 0 or "noto sans arabic" not in family.lower():
        raise DoctorError(
            "Noto Sans Arabic is not available to fontconfig"
        )
    return family


def _find_chromium() -> str:
    for name in _CHROMIUM_CANDIDATES:
        path = shutil.which(name)
        if path:
            return path
    raise DoctorError("headless Chromium is not installed")


def _check_chromium() -> str:
    chromium = _find_chromium()
    result = _run(
        [
            chromium,
            "--headless",
            "--no-sandbox",
            "--disable-gpu",
            "--dump-dom",
            "data:text/html,<p>pilotage-ready</p>",
        ],
        timeout=30,
    )
    if result.returncode != 0 or "pilotage-ready" not in result.stdout:
        raise DoctorError(
            "Chromium did not render a headless page: "
            + _safe_detail(result.stderr)
        )
    version = _run([chromium, "--version"]).stdout.strip()
    return version or chromium


def _check_prepared_environment(config: Any, environment: str) -> str:
    with tempfile.TemporaryDirectory(
        prefix=f"pilotage_doctor_{environment}_"
    ) as directory:
        root = Path(directory)
        context = ToolContext("doctor", config)
        payload = json.loads(
            _execute(
                _SMOKE_CODE[environment],
                environment,
                context,
                workspace=root,
            )
        )
        if payload.get("status") != "success":
            raise DoctorError(
                payload.get("error")
                or payload.get("stderr")
                or "smoke script failed"
            )
        for name in _SMOKE_ARTIFACTS[environment]:
            artifact = root / name
            if not artifact.is_file() or artifact.stat().st_size == 0:
                raise DoctorError(f"smoke test did not create {name}")
        if environment == "pdf":
            if not (root / "bilingual.pdf").read_bytes().startswith(b"%PDF-"):
                raise DoctorError("PDF smoke output has an invalid header")
    return f"{len(_SMOKE_ARTIFACTS[environment])} artifact(s)"


def _check_profile_state(config: Any) -> str:
    state = Path(config.state_dir)
    workspace = Path(config.workspace_dir)
    if not state.is_dir():
        raise DoctorError(f"profile state directory is missing: {state}")
    configured_cwd = config.settings.text("terminal.cwd", "")
    if configured_cwd:
        workspace = Path(configured_cwd).expanduser()
    if not workspace.is_dir():
        raise DoctorError(f"session workspace is missing: {workspace}")
    if not os.access(state, os.R_OK | os.W_OK | os.X_OK):
        raise DoctorError(f"profile state is not accessible: {state}")
    return "state and workspace accessible"


def _check_runtime(config: Any) -> str:
    path = Path(config.state_dir) / ".runtime.lock"
    try:
        record = json.loads(path.read_text(encoding="utf-8"))
        pid = int(record["pid"])
        recorded_state = Path(str(record["state_dir"])).resolve(strict=False)
    except (OSError, ValueError, KeyError, TypeError) as exc:
        raise DoctorError("profile runtime is not running") from exc
    if recorded_state != Path(config.state_dir).resolve(strict=False):
        raise DoctorError("runtime lock belongs to a different profile path")
    try:
        held = runtime_lock_is_held(config.state_dir)
    except RuntimeLockError as exc:
        raise DoctorError(str(exc)) from exc
    if not held:
        raise DoctorError("profile runtime does not own its operating-system lock")
    return f"pid {pid}"


def _check_whatsapp_policy(config: Any) -> str:
    if not config.allowed_senders:
        raise DoctorError("WhatsApp has no explicit allowed senders")
    return "explicit access policy"


def _check_whatsapp_session(config: Any) -> str:
    try:
        validate_whatsapp_session(config.session_dir)
    except WhatsAppSessionError as exc:
        raise DoctorError(str(exc)) from exc
    return "linked device registered"


def _check_whatsapp_bridge(config: Any) -> str:
    path = Path(config.state_dir) / "bridge.pid"
    try:
        record = json.loads(path.read_text(encoding="utf-8"))
        pid = int(record["pid"])
        port = int(record["port"])
        token = str(record["token"])
    except (OSError, ValueError, KeyError, TypeError) as exc:
        raise DoctorError("WhatsApp bridge ownership record is missing") from exc
    if not token or pid <= 0 or not 1 <= port <= 65535:
        raise DoctorError("WhatsApp bridge ownership record is invalid")
    try:
        response = httpx.get(
            f"http://127.0.0.1:{port}/health",
            headers={"x-pilotage-bridge-token": token},
            timeout=3,
        )
        response.raise_for_status()
        payload = response.json()
    except Exception as exc:
        raise DoctorError(
            "WhatsApp bridge health request failed: " + _safe_detail(exc)
        ) from exc
    if int(payload.get("pid", 0)) != pid:
        raise DoctorError("WhatsApp bridge pid does not match its owner record")
    if payload.get("connected") is not True:
        raise DoctorError("WhatsApp bridge is running but not connected")
    return "connected"


def _check_telegram(config: Any) -> str:
    token = os.environ.get("TELEGRAM_BOT_TOKEN", "").strip()
    allowed = os.environ.get("TELEGRAM_ALLOWED_USERS", "").strip()
    if not token or not allowed:
        raise DoctorError("Telegram token or explicit allowed users are missing")
    try:
        response = httpx.get(
            f"https://api.telegram.org/bot{token}/getMe",
            timeout=10,
        )
        response.raise_for_status()
        payload = response.json()
    except Exception as exc:
        raise DoctorError(
            "Telegram getMe failed: " + _safe_detail(exc)
        ) from exc
    if payload.get("ok") is not True:
        raise DoctorError("Telegram rejected the configured bot token")
    return "Bot API accepted the token"


def _check_home_channel(
    whatsapp_config: Any,
    telegram_config: Any,
    *,
    whatsapp_enabled: bool,
    telegram_enabled: bool,
) -> str:
    """Prove the profile has a destination for unsolicited output."""

    configured = []
    if whatsapp_enabled and isinstance(
        getattr(whatsapp_config, "home_origin", None), dict
    ):
        try:
            normalize_whatsapp_chat_id(whatsapp_config.home_origin["chat_id"])
        except (KeyError, TypeError, ValueError) as exc:
            raise DoctorError("WhatsApp home chat is invalid") from exc
        configured.append("WhatsApp")
    if telegram_enabled and isinstance(
        getattr(telegram_config, "home_origin", None), dict
    ):
        try:
            normalize_telegram_home_chat_id(telegram_config.home_origin["chat_id"])
        except (KeyError, TypeError, ValueError) as exc:
            raise DoctorError("Telegram home chat is invalid") from exc
        configured.append("Telegram")
    if not configured:
        raise DoctorError(
            "no enabled channel has a home chat; set WHATSAPP_HOME_CHANNEL or "
            "TELEGRAM_HOME_CHANNEL"
        )
    return "configured for " + " and ".join(configured)


def _sqlcmd_path() -> str:
    configured = os.environ.get("SQLCMD_BIN", "").strip()
    if configured:
        path = Path(configured).expanduser()
        if path.is_file():
            return str(path)
        raise DoctorError(f"SQLCMD_BIN does not exist: {path}")
    found = shutil.which("sqlcmd")
    if found:
        return found
    fallback = Path("/opt/mssql-tools18/bin/sqlcmd")
    if fallback.is_file():
        return str(fallback)
    raise DoctorError("Microsoft sqlcmd is not installed")


def _truthy_env(name: str) -> bool:
    return os.environ.get(name, "").strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }


def _check_sql_connection() -> str:
    missing = [name for name in _SQL_ENV_NAMES if not os.environ.get(name, "").strip()]
    if missing:
        raise DoctorError(
            "missing SQL connection settings: " + ", ".join(missing)
        )
    command = [
        _sqlcmd_path(),
        "-S",
        os.environ["MSSQL_HOST"],
        "-U",
        os.environ["MSSQL_USER"],
        "-d",
        os.environ["MSSQL_DB"],
        "-l",
        "10",
        "-b",
    ]
    if _truthy_env("MSSQL_ENCRYPT"):
        command.append("-N")
    if _truthy_env("MSSQL_TRUST_CERT"):
        command.append("-C")
    command += [
        "-Q",
        "SET NOCOUNT ON; SELECT 1 AS pilotage_ready;",
    ]
    child_env = build_subprocess_env()
    password = child_env.pop("MSSQL_PASSWORD", "")
    child_env["SQLCMDPASSWORD"] = password
    result = _run(command, timeout=20, env=child_env)
    if result.returncode != 0:
        raise DoctorError(
            "sqlcmd connection failed: "
            + _safe_detail(result.stderr or result.stdout)
        )
    return "query completed"


def _check_auth(config: Any) -> str:
    credentials = auth.read_credentials(
        config.credentials_path,
        fallback_path=config.main_credentials_path,
    )
    if not credentials.access_token or not credentials.refresh_token:
        raise DoctorError("Codex OAuth credentials are incomplete")
    return "OAuth credentials readable"


async def _check_model(config: Any) -> str:
    agent = Agent(
        config,
        ConversationStore(path=None),
        disabled_tool_groups=build_registry().groups(),
    )
    try:
        answer = await asyncio.wait_for(
            agent.respond(
                "pilotage-doctor",
                "Reply with the single word OK. This is a deployment readiness probe.",
            ),
            timeout=float(config.request_timeout_seconds) + 10,
        )
    finally:
        await agent.close()
    if not isinstance(answer, str) or not answer.strip():
        raise DoctorError("Codex returned no assistant text")
    return "response received"


def _check_web_configuration() -> str:
    from .tools.web import _get_direct_firecrawl_config

    _get_direct_firecrawl_config()
    return "DDGS and Firecrawl configured"


def _check_stt_configuration() -> str:
    if not (
        os.environ.get("VOICE_TOOLS_OPENAI_KEY", "").strip()
        or os.environ.get("OPENAI_API_KEY", "").strip()
    ):
        raise DoctorError("VOICE_TOOLS_OPENAI_KEY is not configured")
    _check_command("ffmpeg")
    return "OpenAI voice key and ffmpeg configured"


async def collect_report(config: Any, profile_name: str) -> DoctorReport:
    report = DoctorReport(profile=profile_name)
    registry = build_registry()
    groups = set(enabled_groups(config.settings, registry))
    telegram_config = (
        config.for_channel("telegram")
        if callable(getattr(config, "for_channel", None))
        else config
    )
    telegram_groups = set(
        enabled_groups(telegram_config.settings, registry)
    )

    await _probe(report, "Main Python runtime", _check_main_python)
    await _probe(report, "Node.js runtime", _check_node)
    await _probe(
        report,
        "WhatsApp bridge dependencies",
        lambda: _check_bridge_dependencies(config),
    )
    for command in _SYSTEM_COMMANDS:
        await _probe(
            report,
            f"System command: {command}",
            lambda command=command: _check_command(command),
        )
    await _probe(report, "OCR languages", _check_tesseract_languages)
    await _probe(report, "Arabic font", _check_arabic_font)
    await _probe(report, "Headless Chromium", _check_chromium)
    for environment in ("chart", "docs", "excel", "pdf"):
        await _probe(
            report,
            f"Prepared environment: {environment}",
            lambda environment=environment: _check_prepared_environment(
                config, environment
            ),
        )

    await _probe(
        report,
        "Profile state",
        lambda: _check_profile_state(config),
    )
    await _probe(report, "Codex OAuth", lambda: _check_auth(config))
    await _probe_async(report, "Codex model", lambda: _check_model(config))
    await _probe(report, "SQL Server", _check_sql_connection)

    if "web" in groups | telegram_groups:
        await _probe(report, "Web backends", _check_web_configuration)
    if (
        config.settings.flag("stt.enabled", True)
        or telegram_config.settings.flag("stt.enabled", True)
    ):
        await _probe(report, "Speech to text", _check_stt_configuration)

    whatsapp_enabled = config.settings.flag("whatsapp.enabled", False)
    telegram_enabled = telegram_config.settings.flag(
        "telegram.enabled", False
    )
    if whatsapp_enabled or telegram_enabled:
        await _probe(
            report,
            "Home channel",
            lambda: _check_home_channel(
                config,
                telegram_config,
                whatsapp_enabled=whatsapp_enabled,
                telegram_enabled=telegram_enabled,
            ),
        )
    if whatsapp_enabled or telegram_enabled:
        await _probe(
            report,
            "Resident runtime",
            lambda: _check_runtime(config),
        )
    if whatsapp_enabled:
        await _probe(
            report,
            "WhatsApp access policy",
            lambda: _check_whatsapp_policy(config),
        )
        await _probe(
            report,
            "WhatsApp linked session",
            lambda: _check_whatsapp_session(config),
        )
        await _probe(
            report,
            "WhatsApp bridge",
            lambda: _check_whatsapp_bridge(config),
        )
    if telegram_enabled:
        await _probe(
            report,
            "Telegram Bot API",
            lambda: _check_telegram(telegram_config),
        )
    return report


async def run_doctor(
    config: Any,
    profile_name: str,
    *,
    print_fn: Callable[[str], None] = print,
) -> int:
    report = await collect_report(config, profile_name)
    print_fn(f"Pilotage doctor - {profile_name}")
    for check in report.checks:
        marker = "PASS" if check.ok else "FAIL"
        suffix = f" - {check.detail}" if check.detail else ""
        print_fn(f"{marker} {check.name}{suffix}")
    if report.ok:
        print_fn(f"Ready: {len(report.checks)} checks passed.")
        return 0
    print_fn(
        f"Not ready: {len(report.failures)} of {len(report.checks)} checks failed."
    )
    return 1


__all__ = [
    "CheckResult",
    "DoctorReport",
    "collect_report",
    "run_doctor",
]
